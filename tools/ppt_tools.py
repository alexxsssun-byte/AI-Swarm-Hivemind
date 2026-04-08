import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from core.memory import track_file

def generate_slide_deck(project_id: int, deck_type: str, data_summary: dict, project_name: str) -> str:
    """
    Creates a new .pptx pitch deck using python-pptx.
    """
    os.makedirs("output", exist_ok=True)
    filename = f"{project_name.replace(' ', '_')}_{deck_type.replace(' ', '_')}_Deck.pptx"
    filepath = os.path.join("output", filename)
    
    prs = Presentation()
    
    # Slide 1: Title Slide (Layout 0)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"{project_name} - {deck_type}"
    subtitle.text = "Strictly Private and Confidential\nPrepared by AI Investment Banker"
    
    # Slide 2: Executive Summary (Layout 1)
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    tf.text = data_summary.get("exec_summary", "Overview of the transaction and key highlights.")
    
    if "bullets" in data_summary:
        for bullet in data_summary["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1
            
    # Slide 3: Financial Overview & Valuation (Layout 1)
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    slide3.shapes.title.text = "Financial Overview & Valuation"
    tf3 = slide3.shapes.placeholders[1].text_frame
    
    val_text = data_summary.get("valuation_summary", "DCF Yields $X per share. WACC: Y%")
    tf3.text = val_text
    
    # Save
    prs.save(filepath)
    track_file(project_id, filename, "PPTX Presentation")
    
    return filepath

def edit_slide(filepath: str, slide_number: int, edit_instructions: str) -> str:
    """
    Modifies an existing deck.
    """
    if os.path.exists(filepath):
        prs = Presentation(filepath)
        if slide_number <= len(prs.slides):
            slide = prs.slides[slide_number - 1]
            # Just add a text box with the edit for demonstration
            txBox = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = f"[Update]: {edit_instructions}"
            prs.save(filepath)
    return filepath

def get_ppt_tools():
    return [generate_slide_deck, edit_slide]
